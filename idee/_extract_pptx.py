from pptx import Presentation

path = r'g:\AZ\Documents\gestion des appelles telephoniques\idee\presentation_ar.pptx'
out_path = r'g:\AZ\Documents\gestion des appelles telephoniques\idee\_pptx_content.txt'
prs = Presentation(path)

with open(out_path, 'w', encoding='utf-8') as f:
    f.write(f'Nombre de slides: {len(prs.slides)}\n')
    
    for i, slide in enumerate(prs.slides):
        f.write(f'\n{"="*60}\n')
        f.write(f'=== SLIDE {i+1} ===\n')
        f.write(f'{"="*60}\n')
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                text = shape.text.strip()
                f.write(f'  TEXT: {text}\n')
            if shape.has_table:
                table = shape.table
                f.write(f'  TABLE: {len(table.rows)} rows x {len(table.columns)} cols\n')
                for ri, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    f.write(f'    Row {ri}: ' + ' | '.join(cells) + '\n')

print('Done - saved to _pptx_content.txt')
